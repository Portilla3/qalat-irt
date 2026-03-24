#!/usr/bin/env python3
"""
╔══════════════════════════════════════════════════════════════════════════════╗
║   SCRIPT_IRT_Universal_PPTX_Seguimiento.py  —  v1.1                       ║
║   Genera presentación PowerPoint de seguimiento IRT1 → IRT2               ║
║   8 slides · Compatible con cualquier país IRT                             ║
╠══════════════════════════════════════════════════════════════════════════════╣
║  CÓMO USAR:                                                                 ║
║  1. Sube este script + la base Wide IRT                                    ║
║  2. Escribe: "Ejecuta el PPTX Seguimiento IRT"                            ║
║                                                                             ║
║  SLIDES:                                                                    ║
║    1. Portada                                                               ║
║    2. Antecedentes generales (sexo + instrumentos)                        ║
║    3. Sustancia principal IRT1 vs IRT2 (tortas)                           ║
║    4. Días de consumo + Cambio en consumo (barras apiladas)               ║
║    5. Urgencias + Accidentes (2 gráficos lado a lado)                     ║
║    6. Salud psicológica y física + Satisfacción de vida                   ║
║    7. Transgresión total + por tipo                                        ║
║    8. Relaciones interpersonales                                           ║
╚══════════════════════════════════════════════════════════════════════════════╝
"""
import glob, os, unicodedata

def _norm(s):
    return unicodedata.normalize('NFD', str(s).lower()).encode('ascii','ignore').decode()

# ── Detección de país ─────────────────────────────────────────────────────────
_PAISES = {
    'republica_dominicana':'República Dominicana', 'repdomini':'República Dominicana',
    'dominicana':'República Dominicana', 'honduras':'Honduras',
    'panama':'Panamá', 'panam':'Panamá', 'el_salvador':'El Salvador',
    'salvador':'El Salvador', 'mexico':'México', 'mexic':'México',
    'ecuador':'Ecuador', 'peru':'Perú', 'argentina':'Argentina',
    'colombia':'Colombia', 'chile':'Chile', 'bolivia':'Bolivia',
    'paraguay':'Paraguay', 'uruguay':'Uruguay', 'venezuela':'Venezuela',
    'guatemala':'Guatemala', 'costa_rica':'Costa Rica',
    'costarica':'Costa Rica', 'nicaragua':'Nicaragua',
}
def _extraer_pais(filename):
    fn = _norm(str(filename).replace('.','_'))
    for key, nombre in _PAISES.items():
        if key in fn: return nombre
    return None

def _detectar_pais(wide_file):
    import pandas as _pd
    try:
        rs = _pd.read_excel(wide_file, sheet_name='Resumen', header=None)
        for _, row in rs.iterrows():
            for v in row.tolist():
                p = _extraer_pais(str(v))
                if p: return p
    except: pass
    return _extraer_pais(os.path.basename(wide_file))

def auto_archivo_wide():
    candidatos = (
        glob.glob('/mnt/user-data/uploads/*IRT*Wide*.xlsx') +
        glob.glob('/mnt/user-data/uploads/*Wide*IRT*.xlsx') +
        glob.glob('/mnt/user-data/uploads/IRT_Base*.xlsx') +
        glob.glob('/mnt/user-data/outputs/IRT_Base_Wide*.xlsx') +
        glob.glob('/home/claude/IRT_Base_Wide.xlsx'))
    if not candidatos:
        raise FileNotFoundError('⚠  No se encontró la base Wide IRT.')
    uploads = [f for f in candidatos if 'uploads' in f]
    elegido = uploads[0] if uploads else max(candidatos, key=os.path.getsize)
    print(f'  → Base Wide: {os.path.basename(elegido)}')
    return elegido

# ══════════════════════════════════════════════════════════════════════════════
print('=' * 60)
print('  SCRIPT_IRT_Universal_PPTX_Seguimiento  v1.1')
print('=' * 60)

INPUT_FILE  = auto_archivo_wide()
OUTPUT_FILE = '/home/claude/IRT_Presentacion_Seguimiento.pptx'

# ── FILTRO OPCIONAL POR CENTRO ────────────────────────────────────────────────
# Dejar en None para procesar TODOS los centros.
# Ejemplos:
#   FILTRO_CENTRO = None         ← todos los centros
#   FILTRO_CENTRO = "HCHN01"     ← solo ese centro
FILTRO_CENTRO = None
# ─────────────────────────────────────────────────────────────────────────────

import pandas as pd, numpy as np, json, subprocess, sys, warnings
warnings.filterwarnings('ignore')

# ── Carga ─────────────────────────────────────────────────────────────────────
df = pd.read_excel(INPUT_FILE, sheet_name='Base Wide', header=1)
df.columns = [str(c) for c in df.columns]
cols = df.columns.tolist()

_col_centro = next((c for c in cols if any(x in _norm(c) for x in
                    ['codigo del centro', 'servicio de tratamiento',
                     'centro/ servicio', 'codigo centro'])), None)

if FILTRO_CENTRO and _col_centro:
    n_antes = len(df)
    df = df[df[_col_centro].astype(str).str.strip() == FILTRO_CENTRO].copy()
    df = df.reset_index(drop=True)
    print(f'\n⚑  Filtro: "{FILTRO_CENTRO}"  ({n_antes} → {len(df)} pacientes)')
    OUTPUT_FILE = f'/home/claude/IRT_Presentacion_Seguimiento_{FILTRO_CENTRO}.pptx'

N_total = len(df)
mask2   = df['Tiene_IRT2'] == 'Sí'
mask3   = df['Tiene_IRT3'] == 'Sí'
N_irt2  = int(mask2.sum())
N_irt3  = int(mask3.sum())
TIENE_IRT3 = N_irt3 > 0

print(f'\n→ {N_total} pacientes | IRT2: {N_irt2} | IRT3: {N_irt3}')
if N_irt2 == 0:
    raise SystemExit('⚠  Sin IRT2. No se puede generar la presentación.')

# ── País y servicio ────────────────────────────────────────────────────────────
PAIS = _detectar_pais(INPUT_FILE)
if FILTRO_CENTRO:
    NOMBRE_SERVICIO = f'{PAIS}  —  Centro {FILTRO_CENTRO}' if PAIS else f'Centro {FILTRO_CENTRO}'
else:
    NOMBRE_SERVICIO = PAIS if PAIS else 'Servicio de Tratamiento'

# ── Período ───────────────────────────────────────────────────────────────────
MESES = {1:'Enero',2:'Febrero',3:'Marzo',4:'Abril',5:'Mayo',6:'Junio',
         7:'Julio',8:'Agosto',9:'Septiembre',10:'Octubre',11:'Noviembre',12:'Diciembre'}
fecha_col = next((c for c in cols if 'fecha de administracion' in _norm(c)), None)
PERIODO = 'Período no determinado'
if fecha_col:
    fch = pd.to_datetime(df[fecha_col], errors='coerce').dropna()
    anio = pd.Timestamp.now().year
    fch  = fch[(fch.dt.year >= anio-10) & (fch.dt.year <= anio+1)]
    if len(fch):
        f0,f1 = fch.min(), fch.max()
        PERIODO = (f'{MESES[f0.month]} {f0.year}'
                   if f0.year==f1.year and f0.month==f1.month
                   else f'{MESES[f0.month]}–{MESES[f1.month]} {f0.year}'
                   if f0.year==f1.year
                   else f'{MESES[f0.month]} {f0.year} – {MESES[f1.month]} {f1.year}')

print(f'  Servicio: {NOMBRE_SERVICIO} | Período: {PERIODO}')

# ══════════════════════════════════════════════════════════════════════════════
# DETECCIÓN DE COLUMNAS
# ══════════════════════════════════════════════════════════════════════════════
def col_sfx(kws, sfx):
    for c in cols:
        if not c.endswith(sfx): continue
        if all(_norm(k) in _norm(c) for k in kws): return c
    return None

COL_FN   = next((c for c in cols if 'fecha de nacimiento' in _norm(c)), None)
COL_SEXO = next((c for c in cols if _norm(c) in ['sexo','género','genero']), None)

SUST_NOMBRES = {
    'Alcohol':['alcohol'], 'Marihuana':['marihuana','cannabis'],
    'Heroína':['heroina'], 'Cocaína':['cocain'],
    'Fentanilo':['fentanil'], 'Inhalables':['inhalab'],
    'Metanfetamina':['metanfet','cristal'], 'Crack':['crack'],
    'Pasta Base':['pasta base','pasta'], 'Sedantes':['sedant','benzod'],
    'Otra sustancia':['otra sust'],
}
SUST_TOTAL = {}
for sust, kws in SUST_NOMBRES.items():
    entry = {}
    for sfx in ['_IRT1','_IRT2','_IRT3']:
        for c in cols:
            if not c.endswith(sfx): continue
            nc = _norm(c)
            if any(_norm(k) in nc for k in kws) and ('total' in nc or '(0-28)' in nc):
                entry[sfx] = c; break
    if entry: SUST_TOTAL[sust] = entry
SUST_ACTIVAS = list(SUST_TOTAL.keys())

COL_SP   = {sfx: col_sfx(['sustancia','principal'], sfx) for sfx in ['_IRT1','_IRT2','_IRT3']}
COL_SPSI = {sfx: col_sfx(['salud','psicol'], sfx)         for sfx in ['_IRT1','_IRT2','_IRT3']}
COL_SFIS = {sfx: col_sfx(['salud','fis'],    sfx)         for sfx in ['_IRT1','_IRT2','_IRT3']}
COL_URG  = {sfx: next((c for c in cols if c.endswith(sfx) and '5)' in c and
                        any(k in c.lower() for k in ['urgencia','hospitali','emergencia'])), None)
            for sfx in ['_IRT1','_IRT2','_IRT3']}
COL_ACC  = {sfx: next((c for c in cols if c.endswith(sfx) and '6)' in c and
                        'accidente' in c.lower()), None)
            for sfx in ['_IRT1','_IRT2','_IRT3']}
TRANS_DEF = {'Robo / Hurto':'robo','Venta de sustancias':'venta',
             'Violencia':'violencia','VIF':'intraf','Detenido':'detenido'}
TRANS_COLS = {}
for nombre, kw in TRANS_DEF.items():
    entry = {sfx: next((c for c in cols if c.endswith(sfx) and kw in c.lower()), None)
             for sfx in ['_IRT1','_IRT2']}
    if any(entry.values()): TRANS_COLS[nombre] = entry

REL_MAP = {'Padre':['padre'],'Madre':['madre'],'Hijos':['hijos','hijo'],
           'Hermanos':['hermanos'],'Pareja':['pareja'],'Amigos':['amigos'],'Otros':['otros']}
REL_COLS = {}
for vin, kws in REL_MAP.items():
    for sfx in ['_IRT1','_IRT2']:
        for c in cols:
            if not c.endswith(sfx): continue
            nc = _norm(c)
            if '14)' not in c and 'relaci' not in nc: continue
            if any(_norm(k) in nc for k in kws):
                REL_COLS.setdefault(vin,{})[sfx]=c; break

SAT_MAP = {
    'Vida general':[['16)'],['satisfac','vida']],
    'Lugar donde vive':[['17)'],['satisfac','lugar']],
    'Situación laboral':[['18)'],['satisfac','labor','educac']],
    'Tiempo libre':[['19)'],['satisfac','tiempo']],
    'Cap. económica':[['20)'],['satisfac','econom']],
}
SAT_COLS = {}
for dim,(nums,kws) in SAT_MAP.items():
    for sfx in ['_IRT1','_IRT2']:
        for c in cols:
            if not c.endswith(sfx): continue
            nc = _norm(c)
            if any(n in c for n in nums) and any(_norm(k) in nc for k in kws):
                SAT_COLS.setdefault(dim,{})[sfx]=c; break

# ══════════════════════════════════════════════════════════════════════════════
# CALCULAR DATOS
# ══════════════════════════════════════════════════════════════════════════════
print('\n→ Calculando datos...')

def norm_sust(s):
    if pd.isna(s) or str(s).strip() in ['0','']: return None
    s = _norm(str(s))
    if any(x in s for x in ['alcohol','cerveza','licor']): return 'Alcohol'
    if any(x in s for x in ['marihu','cannabis','marij']): return 'Marihuana'
    if any(x in s for x in ['crack','piedra','paco']):     return 'Crack'
    if any(x in s for x in ['pasta base','pasta']):        return 'Pasta Base'
    if any(x in s for x in ['cocain','perico']):           return 'Cocaína'
    if any(x in s for x in ['fentanil']):                  return 'Fentanilo'
    if any(x in s for x in ['inhalab']):                   return 'Inhalables'
    if any(x in s for x in ['metanfet','cristal']):        return 'Metanfetamina'
    if any(x in s for x in ['sedant','benzod']):           return 'Sedantes'
    if any(x in s for x in ['heroina','opiod','morfin']):  return 'Heroína'
    if any(x in s for x in ['tabaco','nicot']):            return 'Tabaco'
    return 'Otra sustancia'

def prom(col, msk):
    if col is None: return None
    v = pd.to_numeric(df.loc[msk,col], errors='coerce').dropna()
    return round(float(v.mean()),1) if len(v) else None

def sino_pct(col, msk):
    if col is None: return None, None
    v = df.loc[msk,col].dropna().astype(str).str.strip().str.lower()
    nv = len(v)
    nsi = int(v.isin(['sí','si']).sum())
    return nsi, round(nsi/nv*100,1) if nv else 0.0

hoy = pd.Timestamp.now()
pct_seg = round(N_irt2/N_total*100,1)

# Sexo
sexo_data = {}
if COL_SEXO:
    sc = df.loc[mask2, COL_SEXO].value_counts(dropna=True)
    sexo_data = {k:int(v) for k,v in sc.items()}

# Edad
edad_data = {}
if COL_FN:
    edades = ((hoy - pd.to_datetime(df.loc[mask2,COL_FN], errors='coerce')).dt.days/365.25).dropna()
    edad_data = {'mean': round(edades.mean(),1), 'min': int(edades.min()),
                 'max': int(edades.max())} if len(edades) else {}

# Sustancia principal IRT1 vs IRT2
def dist_sp(sfx, msk):
    c = COL_SP.get(sfx)
    if not c: return {}
    return df.loc[msk,c].apply(norm_sust).dropna().value_counts().to_dict()

sp1 = dist_sp('_IRT1', mask2)
sp2 = dist_sp('_IRT2', mask2)

# Días consumo por sustancia IRT1 vs IRT2
dias_data = []
for sust in SUST_ACTIVAS:
    entry = SUST_TOTAL[sust]
    c1 = entry.get('_IRT1'); c2 = entry.get('_IRT2')
    if not c1: continue
    v1c = pd.to_numeric(df.loc[mask2,c1],errors='coerce'); v1c = v1c[v1c>0].dropna()
    v2c = pd.to_numeric(df.loc[mask2,c2],errors='coerce') if c2 else pd.Series(dtype=float)
    v2c = v2c[v2c>0].dropna()
    m1 = round(float(v1c.mean()),1) if len(v1c) else 0.0
    m2 = round(float(v2c.mean()),1) if len(v2c) else 0.0
    if m1 > 0 or m2 > 0:
        dias_data.append({'label':sust,'irt1':m1,'irt2':m2})
dias_data.sort(key=lambda x:-x['irt1'])

# Cambio en consumo (barras apiladas) — solo consumidores IRT1
cambio_data = []
for sust in SUST_ACTIVAS:
    entry = SUST_TOTAL[sust]
    c1 = entry.get('_IRT1'); c2 = entry.get('_IRT2')
    if not c1 or not c2: continue
    sp1col = COL_SP.get('_IRT1')
    if not sp1col: continue
    # Solo personas cuya sustancia principal es esta
    msk_pp = mask2 & (df[sp1col].apply(norm_sust) == sust)
    if not msk_pp.any(): continue
    v1 = pd.to_numeric(df.loc[msk_pp,c1],errors='coerce')
    v2 = pd.to_numeric(df.loc[msk_pp,c2],errors='coerce')
    ok = v1.notna() & v2.notna()
    # También incluir cualquier consumidor en IRT1
    msk_cons = mask2 & (pd.to_numeric(df[c1],errors='coerce').fillna(0) > 0) & pd.to_numeric(df[c2],errors='coerce').notna()
    v1b = pd.to_numeric(df.loc[msk_cons,c1],errors='coerce')
    v2b = pd.to_numeric(df.loc[msk_cons,c2],errors='coerce')
    ok2 = v1b.notna() & v2b.notna()
    n_ok = int(ok2.sum())
    if n_ok == 0: continue
    n_abs = int((v2b[ok2]==0).sum())
    n_dis = int(((v2b[ok2]>0)&(v2b[ok2]<v1b[ok2])).sum())
    n_sc  = int((v2b[ok2]==v1b[ok2]).sum())
    n_emp = int((v2b[ok2]>v1b[ok2]).sum())
    pct = lambda n: round(n/n_ok*100,1)
    cambio_data.append({'label':sust,'n':n_ok,
        'abs':pct(n_abs),'dis':pct(n_dis),'sin':pct(n_sc),'emp':pct(n_emp),
        'combo':round((n_abs+n_dis)/n_ok*100,1)})

# Urgencias y accidentes
n_urg1, pct_urg1 = sino_pct(COL_URG.get('_IRT1'), mask2)
n_urg2, pct_urg2 = sino_pct(COL_URG.get('_IRT2'), mask2)
n_acc1, pct_acc1 = sino_pct(COL_ACC.get('_IRT1'), mask2)
n_acc2, pct_acc2 = sino_pct(COL_ACC.get('_IRT2'), mask2)

# Salud
salud_data = []
for nombre, cd in [('Salud Psicológica (0–10)', COL_SPSI),
                   ('Salud Física (0–10)',       COL_SFIS)]:
    m1 = prom(cd.get('_IRT1'), mask2)
    m2 = prom(cd.get('_IRT2'), mask2)
    if m1 is not None or m2 is not None:
        salud_data.append({'label':nombre,'irt1':m1 or 0,'irt2':m2 or 0})

# Satisfacción
sat_data = []
for dim, cd in SAT_COLS.items():
    m1 = prom(cd.get('_IRT1'), mask2)
    m2 = prom(cd.get('_IRT2'), mask2)
    if m1 is not None or m2 is not None:
        sat_data.append({'label':dim,'irt1':m1 or 0,'irt2':m2 or 0})

# Transgresión
def mask_trans(sfx, msk):
    cols_tr = [d.get(sfx) for d in TRANS_COLS.values() if d.get(sfx)]
    if not cols_tr: return pd.Series(False, index=df.index[msk])
    return pd.concat([pd.to_numeric(df.loc[msk,c],errors='coerce').fillna(0)>0
                      for c in cols_tr], axis=1).any(axis=1)

n_tr1 = int(mask_trans('_IRT1',mask2).sum())
n_tr2 = int(mask_trans('_IRT2',mask2).sum())
pct_tr1 = round(n_tr1/N_irt2*100,1); pct_tr2 = round(n_tr2/N_irt2*100,1)

tipos_tr = []
for tipo, cd in TRANS_COLS.items():
    c1=cd.get('_IRT1'); c2=cd.get('_IRT2')
    v1 = pd.to_numeric(df.loc[mask2,c1],errors='coerce').fillna(0) if c1 else pd.Series([0]*N_irt2)
    v2 = pd.to_numeric(df.loc[mask2,c2],errors='coerce').fillna(0) if c2 else pd.Series([0]*N_irt2)
    p1 = round((v1>0).sum()/N_irt2*100,1); p2 = round((v2>0).sum()/N_irt2*100,1)
    if p1>0 or p2>0:
        tipos_tr.append({'label':tipo,'irt1':p1,'irt2':p2})

# Relaciones
CATS_REL = ['Excelente','Buena','Ni buena ni mala','Mala','Muy mala']
rel_data = []
for vin, cd in REL_COLS.items():
    row_rel = {'label': vin}
    for sfx, lbl in [('_IRT1','irt1'),('_IRT2','irt2')]:
        c = cd.get(sfx)
        if not c: row_rel[f'pos_{lbl}']=0; row_rel[f'neg_{lbl}']=0; continue
        vals = df.loc[mask2,c].dropna()
        vals = vals[vals.astype(str).str.lower()!='no aplica']
        nv = len(vals)
        if nv == 0: row_rel[f'pos_{lbl}']=0; row_rel[f'neg_{lbl}']=0; continue
        pos = int((vals.astype(str).isin(['Excelente','Buena'])).sum())
        neg = int((vals.astype(str).isin(['Mala','Muy mala'])).sum())
        row_rel[f'pos_{lbl}'] = round(pos/nv*100,1)
        row_rel[f'neg_{lbl}'] = round(neg/nv*100,1)
    rel_data.append(row_rel)

print('  ✓ Datos calculados')
print(f'  Sustancias: {SUST_ACTIVAS}')
print(f'  Cambio: {len(cambio_data)} sustancias | Trans. IRT1: {pct_tr1}% → IRT2: {pct_tr2}%')

# ══════════════════════════════════════════════════════════════════════════════
# JSON INTERMEDIO
# ══════════════════════════════════════════════════════════════════════════════
data = {
    'meta': {
        'servicio': NOMBRE_SERVICIO,
        'periodo':  PERIODO,
        'N_total':  N_total,
        'N_irt2':   N_irt2,
        'N_irt3':   N_irt3,
        'pct_seg':  pct_seg,
    },
    'sexo':    [{'label':k,'n':v,'pct':round(v/N_irt2*100,1)} for k,v in sexo_data.items()],
    'edad':    edad_data,
    'sp1':     [{'label':k,'n':v,'pct':round(v/N_irt2*100,1)} for k,v in sp1.items()],
    'sp2':     [{'label':k,'n':v,'pct':round(v/N_irt2*100,1)} for k,v in sp2.items()],
    'dias':    dias_data,
    'cambio':  cambio_data,
    'urgencias': {
        'n1':   n_urg1, 'pct1': pct_urg1,
        'n2':   n_urg2, 'pct2': pct_urg2,
    },
    'accidentes': {
        'n1':   n_acc1, 'pct1': pct_acc1,
        'n2':   n_acc2, 'pct2': pct_acc2,
    },
    'salud':   salud_data,
    'sat':     sat_data,
    'transTotal': {'irt1': pct_tr1, 'irt2': pct_tr2, 'n1': n_tr1, 'n2': n_tr2},
    'transTipos': tipos_tr,
    'rel':     rel_data,
}


# ══════════════════════════════════════════════════════════════════════════════
# PYTHON-PPTX — Construcción del PowerPoint (sin Node.js)
# ══════════════════════════════════════════════════════════════════════════════
import io, matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

C_DARK = RGBColor(0x1F,0x38,0x64); C_MID  = RGBColor(0x2E,0x75,0xB6)
C_LIGHT= RGBColor(0xBD,0xD7,0xEE); C_ACC  = RGBColor(0x00,0xB0,0xF0)
C_WHITE= RGBColor(0xFF,0xFF,0xFF); C_GRAY = RGBColor(0x59,0x59,0x59)
MC_I1='#1F3864'; MC_I2='#00B0F0'
MC_ABS='#1F3864'; MC_DIS='#375623'; MC_SC='#BFBFBF'; MC_EMP='#C00000'
PIE_COLS=['#2E75B6','#1F3864','#00B0F0','#9DC3E6','#70AD47',
          '#4472C4','#D9D9D9','#C00000','#ED7D31','#FFC000']

SLIDE_W=Inches(10); SLIDE_H=Inches(5.625)

def add_rect(sl,x,y,w,h,fill):
    s=sl.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb=fill; s.line.fill.background()

def add_txt(sl,text,x,y,w,h,size=11,bold=False,color=None,align=PP_ALIGN.LEFT,italic=False):
    tb=sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=str(text)
    r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic
    if color: r.font.color.rgb=color

def hdr(sl,txt):
    add_rect(sl,0,0,10,0.72,C_DARK); add_rect(sl,5.5,0,4.5,0.72,C_ACC)
    add_txt(sl,txt,0.25,0.05,9.5,0.62,size=18,bold=True,color=C_WHITE)

def ftr(sl,txt):
    add_txt(sl,txt,0.25,5.32,9.5,0.25,size=8,color=C_GRAY,align=PP_ALIGN.CENTER,italic=True)

def divv(sl,x):
    ln=sl.shapes.add_shape(1,Inches(x),Inches(0.78),Inches(0.02),Inches(4.85))
    ln.fill.solid(); ln.fill.fore_color.rgb=RGBColor(0xD9,0xD9,0xD9); ln.line.fill.background()

def fig2img(sl,fig,x,y,w,h):
    buf=io.BytesIO()
    fig.savefig(buf,format='png',dpi=130,bbox_inches='tight',facecolor='white')
    buf.seek(0); plt.close(fig)
    sl.shapes.add_picture(buf,Inches(x),Inches(y),Inches(w),Inches(h))

def axstyle(ax,horiz=False):
    ax.set_facecolor('white')
    (ax.xaxis if horiz else ax.yaxis).grid(True,color='#E2E8F0',linewidth=0.6,zorder=0)
    ax.set_axisbelow(True)
    for sp in ['top','right']: ax.spines[sp].set_visible(False)
    ax.spines['left'].set_color('#D0D0D0'); ax.spines['bottom'].set_color('#D0D0D0')

def g_pie(labels,values):
    fig,ax=plt.subplots(figsize=(4.2,3.5))
    wedges,_,at=ax.pie(values,labels=None,colors=PIE_COLS[:len(values)],
        autopct=lambda p:f'{p:.1f}%' if p>4 else '',startangle=140,pctdistance=0.72,
        wedgeprops={'edgecolor':'white','linewidth':1.5})
    for a in at: a.set_fontsize(9); a.set_color('white'); a.set_fontweight('bold')
    ax.legend(wedges,[f'{l} (n={v})' for l,v in zip(labels,values)],
              loc='lower center',bbox_to_anchor=(0.5,-0.18),ncol=2,fontsize=8,frameon=False)
    ax.set_aspect('equal'); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

def g_barras_dobles(labels,vals1,vals2,lab1='IRT1',lab2='IRT2',ylabel='',max_val=None):
    x=np.arange(len(labels)); ww=0.35
    fig,ax=plt.subplots(figsize=(max(4.5,len(labels)*0.85),3.5))
    b1=ax.bar(x-ww/2,vals1,ww,color=MC_I1,label=lab1,zorder=3)
    b2=ax.bar(x+ww/2,vals2,ww,color=MC_I2,label=lab2,zorder=3)
    for b,v in zip(list(b1)+list(b2),vals1+vals2):
        if v>0: ax.text(b.get_x()+b.get_width()/2,b.get_height()+0.3,
                        str(v),ha='center',va='bottom',fontsize=8,fontweight='bold',color='#333')
    ax.set_xticks(x); ax.set_xticklabels(labels,fontsize=9,rotation=20,ha='right')
    if ylabel: ax.set_ylabel(ylabel,fontsize=8,color='#595959')
    if max_val: ax.set_ylim(0,max_val*1.15)
    ax.legend(fontsize=8,frameon=False)
    axstyle(ax); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

def g_apilado(cambio_data):
    if not cambio_data: return None
    labs=[c['label'] for c in cambio_data]
    abs_=[c['abs'] for c in cambio_data]; dis_=[c['dis'] for c in cambio_data]
    sc_ =[c['sin'] for c in cambio_data]; emp_=[c['emp'] for c in cambio_data]
    x=np.arange(len(labs))
    fig,ax=plt.subplots(figsize=(max(4.5,len(labs)*0.85),3.5))
    ax.bar(x,abs_,color=MC_ABS,label='Abstinencia',zorder=3)
    ax.bar(x,dis_,bottom=abs_,color=MC_DIS,label='Disminuyó',zorder=3)
    ax.bar(x,sc_,bottom=[a+d for a,d in zip(abs_,dis_)],color=MC_SC,label='Sin cambio',zorder=3)
    ax.bar(x,emp_,bottom=[a+d2+s for a,d2,s in zip(abs_,dis_,sc_)],color=MC_EMP,label='Empeoró',zorder=3)
    for i,(a,d2,s,e) in enumerate(zip(abs_,dis_,sc_,emp_)):
        y_=0
        for val,col in [(a,'white'),(d2,'white'),(s,'#333'),(e,'white')]:
            if val>9: ax.text(i,y_+val/2,f'{val:.0f}%',ha='center',va='center',
                              fontsize=7.5,color=col,fontweight='bold')
            y_+=val
    ax.set_xticks(x); ax.set_xticklabels(labs,fontsize=8,rotation=20,ha='right')
    ax.set_ylim(0,115); ax.set_ylabel('% consumidores al ingreso',fontsize=8,color='#595959')
    ax.legend(fontsize=7.5,frameon=False,ncol=2,loc='upper right')
    axstyle(ax); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

def g_salud(salud_data):
    if not salud_data: return None
    labs=[s['label'] for s in salud_data]
    v1=[s['irt1'] for s in salud_data]; v2=[s['irt2'] for s in salud_data]
    y=np.arange(len(labs)); ww=0.35
    fig,ax=plt.subplots(figsize=(4.5,2.5))
    b1=ax.barh(y-ww/2,v1,ww,color=MC_I1,label='Ingreso (IRT1)',zorder=3)
    b2=ax.barh(y+ww/2,v2,ww,color=MC_I2,label='Seguimiento (IRT2)',zorder=3)
    for b,v in zip(list(b1)+list(b2),v1+v2):
        ax.text(b.get_width()+0.1,b.get_y()+b.get_height()/2,
                str(v),va='center',fontsize=9,fontweight='bold',color='#333')
    ax.set_yticks(y); ax.set_yticklabels(labs,fontsize=9)
    ax.set_xlim(0,12); ax.axvline(x=5,color='#BFBFBF',linestyle='--',linewidth=0.8)
    ax.legend(fontsize=8,frameon=False,loc='lower right')
    axstyle(ax,horiz=True); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

# ── Construcción slides ───────────────────────────────────────────────────────
prs=Presentation()
prs.slide_width=SLIDE_W; prs.slide_height=SLIDE_H
blank=prs.slide_layouts[6]
TITULO=f'Seguimiento IRT1 vs IRT2 · {NOMBRE_SERVICIO}'
PIE_TXT=f'N seguimiento = {N_irt2}  ·  {NOMBRE_SERVICIO}  ·  {PERIODO}'
pct_seg_val=round(N_irt2/N_total*100,1) if N_total else 0

# SLIDE 1: PORTADA
sl=prs.slides.add_slide(blank)
add_rect(sl,0,0,4.0,5.625,C_DARK); add_rect(sl,3.1,0,1.5,5.625,C_ACC)
add_txt(sl,'Seguimiento',0.25,1.6,3.2,0.7,size=22,bold=True,color=C_WHITE)
add_txt(sl,'IRT1 vs IRT2',0.25,2.35,3.2,0.55,size=12,color=C_LIGHT)
add_txt(sl,NOMBRE_SERVICIO.upper(),0.25,3.1,3.2,0.6,size=13,bold=True,color=C_WHITE)
add_txt(sl,PERIODO,0.25,3.75,3.2,0.4,size=11,color=C_LIGHT)
add_txt(sl,f'N ingreso: {N_total}  ·  Con IRT2: {N_irt2} ({pct_seg_val}%)',0.25,4.3,3.2,0.4,size=10,color=C_LIGHT)
add_txt(sl,'IRT 1 → IRT 2\nSeguimiento',4.6,1.55,5.1,1.4,size=26,bold=True,color=C_GRAY,align=PP_ALIGN.CENTER)
add_txt(sl,NOMBRE_SERVICIO.upper(),4.6,3.1,5.1,0.45,size=15,bold=True,color=C_MID,align=PP_ALIGN.CENTER)
add_txt(sl,PERIODO,4.6,3.62,5.1,0.35,size=13,bold=True,color=C_MID,align=PP_ALIGN.CENTER)
add_txt(sl,f'N ingreso: {N_total}  ·  Con IRT2: {N_irt2} ({pct_seg_val}%)',4.6,4.1,5.1,0.35,size=10,color=C_GRAY,align=PP_ALIGN.CENTER)

# SLIDE 2: ANTECEDENTES GENERALES (Sexo + Instrumentos)
sl=prs.slides.add_slide(blank)
hdr(sl,TITULO); divv(sl,4.95)
add_txt(sl,'DISTRIBUCIÓN POR SEXO\n(participantes con IRT2)',0.25,0.82,4.5,0.55,size=12,bold=True,color=C_GRAY)
if sexo_data:
    labs_s=list(sexo_data.keys()); vals_s=list(sexo_data.values())
    fig_s,ax=plt.subplots(figsize=(4.2,3.5))
    bars=ax.bar(labs_s,vals_s,color=['#2E75B6','#9DC3E6','#BDD7EE'][:len(labs_s)],width=0.5,zorder=3)
    for b,v in zip(bars,vals_s):
        ax.text(b.get_x()+b.get_width()/2,b.get_height()+0.3,
                f'{v}\n({round(v/N_irt2*100,1)}%)',ha='center',va='bottom',fontsize=10,fontweight='bold',color='#333')
    ax.set_ylim(0,max(vals_s)*1.4 if vals_s else 1)
    axstyle(ax); fig_s.patch.set_facecolor('white'); fig_s.tight_layout()
    fig2img(sl,fig_s,0.25,1.45,4.5,3.85)
# Derecha: instrumentos completados
add_txt(sl,'INSTRUMENTOS COMPLETADOS',5.15,0.82,4.6,0.35,size=12,bold=True,color=C_GRAY)
for i,(lbl,n_v,pct_v) in enumerate([
    ('IRT1 — Ingreso',N_total,100),
    ('IRT2 — Seguimiento 3m',N_irt2,pct_seg_val),
    ('IRT3 — Seguimiento 6m',N_irt3,round(N_irt3/N_total*100,1) if N_total else 0)]):
    y_box=1.35+i*1.25
    add_rect(sl,5.2,y_box,4.5,1.0,RGBColor(0xEE,0xF4,0xFB))
    add_txt(sl,lbl,5.3,y_box+0.05,4.3,0.3,size=10,bold=True,color=C_DARK)
    add_txt(sl,f'n = {n_v}  ({pct_v}%)',5.3,y_box+0.4,4.3,0.35,size=14,bold=True,color=C_MID)
ftr(sl,PIE_TXT)

# SLIDE 3: SUSTANCIA PRINCIPAL IRT1 vs IRT2
sl=prs.slides.add_slide(blank)
hdr(sl,TITULO); divv(sl,4.95)
add_txt(sl,'SUSTANCIA PRINCIPAL\nIngreso (IRT1)',0.25,0.82,4.5,0.55,size=12,bold=True,color=C_GRAY)
if sp1:
    labs1=list(sp1.keys()); vals1=list(sp1.values())
    fig_sp1=g_pie(labs1,vals1)
    fig2img(sl,fig_sp1,0.2,1.45,4.5,3.9)
add_txt(sl,'SUSTANCIA PRINCIPAL\nSeguimiento (IRT2)',5.15,0.82,4.6,0.55,size=12,bold=True,color=C_GRAY)
if sp2:
    labs2=list(sp2.keys()); vals2=list(sp2.values())
    fig_sp2=g_pie(labs2,vals2)
    fig2img(sl,fig_sp2,5.1,1.45,4.65,3.9)
ftr(sl,PIE_TXT)

# SLIDE 4: DÍAS CONSUMO + CAMBIO EN CONSUMO
sl=prs.slides.add_slide(blank)
hdr(sl,TITULO); divv(sl,4.95)
add_txt(sl,'DÍAS DE CONSUMO\nIRT1 vs IRT2 (promedio)',0.25,0.82,4.5,0.55,size=11,bold=True,color=C_MID)
if dias_data:
    labs_d=[d['label'] for d in dias_data]
    v1_d=[d['irt1'] for d in dias_data]; v2_d=[d['irt2'] for d in dias_data]
    fig_d=g_barras_dobles(labs_d,v1_d,v2_d,ylabel='Días (0–28)',max_val=28)
    fig2img(sl,fig_d,0.2,1.45,4.5,3.85)
add_txt(sl,'CAMBIO EN EL CONSUMO\n% de consumidores al ingreso',5.15,0.82,4.6,0.55,size=11,bold=True,color=C_MID)
fig_cb=g_apilado(cambio_data)
if fig_cb: fig2img(sl,fig_cb,5.15,1.45,4.6,3.85)
ftr(sl,PIE_TXT)

# SLIDE 5: URGENCIAS + ACCIDENTES
sl=prs.slides.add_slide(blank)
hdr(sl,TITULO); divv(sl,4.95)
add_txt(sl,'URGENCIAS U HOSPITALIZACIÓN\nPor consumo de sustancias',0.25,0.82,4.5,0.55,size=12,bold=True,color=C_GRAY)
n1_u=data['urgencias']['n1']; p1_u=data['urgencias']['pct1']
n2_u=data['urgencias']['n2']; p2_u=data['urgencias']['pct2']
fig_u,ax=plt.subplots(figsize=(3.8,3.2))
ax.bar(['IRT1','IRT2'],[p1_u,p2_u],color=[MC_I1,MC_I2],width=0.5,zorder=3)
for x_,v in enumerate([p1_u,p2_u]):
    ax.text(x_,v+0.5,f'{v}%',ha='center',va='bottom',fontsize=12,fontweight='bold',color='#333')
ax.set_ylim(0,max(p1_u,p2_u)*1.5 if max(p1_u,p2_u)>0 else 5)
ax.set_ylabel('% personas',fontsize=8,color='#595959')
axstyle(ax); fig_u.patch.set_facecolor('white'); fig_u.tight_layout()
fig2img(sl,fig_u,0.4,1.45,4.0,3.85)

add_txt(sl,'ACCIDENTES RELACIONADOS\nCon el consumo',5.15,0.82,4.6,0.55,size=12,bold=True,color=C_GRAY)
n1_a=data['accidentes']['n1']; p1_a=data['accidentes']['pct1']
n2_a=data['accidentes']['n2']; p2_a=data['accidentes']['pct2']
fig_a,ax=plt.subplots(figsize=(3.8,3.2))
ax.bar(['IRT1','IRT2'],[p1_a,p2_a],color=[MC_I1,MC_I2],width=0.5,zorder=3)
for x_,v in enumerate([p1_a,p2_a]):
    ax.text(x_,v+0.5,f'{v}%',ha='center',va='bottom',fontsize=12,fontweight='bold',color='#333')
ax.set_ylim(0,max(p1_a,p2_a)*1.5 if max(p1_a,p2_a)>0 else 5)
ax.set_ylabel('% personas',fontsize=8,color='#595959')
axstyle(ax); fig_a.patch.set_facecolor('white'); fig_a.tight_layout()
fig2img(sl,fig_a,5.35,1.45,4.0,3.85)
ftr(sl,PIE_TXT)

# SLIDE 6: SALUD + SATISFACCIÓN
sl=prs.slides.add_slide(blank)
hdr(sl,TITULO); divv(sl,4.95)
add_txt(sl,'AUTOPERCEPCIÓN DE SALUD (0–10)\nIRT1 vs IRT2',0.25,0.82,4.5,0.55,size=11,bold=True,color=C_MID)
fig_sal=g_salud(salud_data)
if fig_sal: fig2img(sl,fig_sal,0.2,1.45,4.5,3.85)
add_txt(sl,'SATISFACCIÓN DE VIDA (0–10)\nIRT1 vs IRT2',5.15,0.82,4.6,0.55,size=11,bold=True,color=C_MID)
if sat_data:
    labs_sat=[s['label'] for s in sat_data]
    v1_s=[s['irt1'] for s in sat_data]; v2_s=[s['irt2'] for s in sat_data]
    fig_sat=g_barras_dobles(labs_sat,v1_s,v2_s,ylabel='Promedio (0–10)',max_val=10)
    fig2img(sl,fig_sat,5.15,1.45,4.6,3.85)
ftr(sl,PIE_TXT)

# SLIDE 7: TRANSGRESIÓN
sl=prs.slides.add_slide(blank)
hdr(sl,TITULO); divv(sl,4.95)
n_t1=data['transTotal']['n1']; p_t1=data['transTotal']['irt1']
n_t2=data['transTotal']['n2']; p_t2=data['transTotal']['irt2']
add_txt(sl,'TRANSGRESIÓN A LA NORMA SOCIAL\nIRT1 vs IRT2',0.25,0.82,4.5,0.55,size=12,bold=True,color=C_GRAY)
fig_tr,ax=plt.subplots(figsize=(3.8,3.2))
ax.bar(['IRT1','IRT2'],[p_t1,p_t2],color=[MC_I1,MC_I2],width=0.5,zorder=3)
for x_,v in enumerate([p_t1,p_t2]):
    ax.text(x_,v+0.5,f'{v}%',ha='center',va='bottom',fontsize=12,fontweight='bold',color='#333')
ax.set_ylim(0,max(p_t1,p_t2)*1.5 if max(p_t1,p_t2)>0 else 5)
ax.set_ylabel('% personas',fontsize=8,color='#595959')
axstyle(ax); fig_tr.patch.set_facecolor('white'); fig_tr.tight_layout()
fig2img(sl,fig_tr,0.4,1.45,4.0,3.85)

add_txt(sl,'DISTRIBUCIÓN POR TIPO\nIRT1 vs IRT2',5.15,0.82,4.6,0.55,size=11,bold=True,color=C_GRAY)
if tipos_tr:
    labs_tt=[t['label'] for t in tipos_tr]
    v1_t=[t['irt1'] for t in tipos_tr]; v2_t=[t['irt2'] for t in tipos_tr]
    fig_tt=g_barras_dobles(labs_tt,v1_t,v2_t,ylabel='%')
    fig2img(sl,fig_tt,5.1,1.45,4.65,3.85)
ftr(sl,PIE_TXT)

# SLIDE 8: RELACIONES INTERPERSONALES
sl=prs.slides.add_slide(blank)
hdr(sl,TITULO)
add_txt(sl,'RELACIONES INTERPERSONALES\n% positivas (Excelente + Buena) · IRT1 vs IRT2',0.25,0.82,9.5,0.55,size=13,bold=True,color=C_MID)
if rel_data:
    labs_r=[r['label'] for r in rel_data]
    v1_r=[r['pos_irt1'] for r in rel_data]; v2_r=[r['pos_irt2'] for r in rel_data]
    y_r=np.arange(len(labs_r)); ww=0.35
    fig_r,ax=plt.subplots(figsize=(7,max(3,len(labs_r)*0.6)))
    b1=ax.barh(y_r-ww/2,v1_r,ww,color=MC_I1,label='Ingreso (IRT1)',zorder=3)
    b2=ax.barh(y_r+ww/2,v2_r,ww,color=MC_I2,label='Seguimiento (IRT2)',zorder=3)
    for b,v in zip(list(b1)+list(b2),v1_r+v2_r):
        ax.text(b.get_width()+0.5,b.get_y()+b.get_height()/2,
                f'{v}%',va='center',fontsize=9,fontweight='bold',color='#333')
    ax.set_yticks(y_r); ax.set_yticklabels(labs_r,fontsize=10)
    ax.set_xlim(0,120); ax.set_xlabel('% positivo',fontsize=9,color='#595959')
    ax.legend(fontsize=9,frameon=False,loc='lower right')
    axstyle(ax,horiz=True); fig_r.patch.set_facecolor('white'); fig_r.tight_layout()
    fig2img(sl,fig_r,1.2,1.45,7.6,3.85)
ftr(sl,PIE_TXT)

prs.save(OUTPUT_FILE)
print(f'\n{"="*60}')
print(f'  ✅  LISTO  →  {OUTPUT_FILE}')
print(f'      {N_irt2}/{N_total} pacientes con IRT2  ·  {PERIODO}')
print(f'{"="*60}')
