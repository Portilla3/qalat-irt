#!/usr/bin/env python3
"""
╔══════════════════════════════════════════════════════════════════════════════╗
║   SCRIPT_IRT_Universal_PPTX_Caracterizacion.py  —  v1.1                   ║
║   Genera presentación PowerPoint de caracterización al ingreso (IRT1)    ║
║   8 slides · Compatible con cualquier país IRT                            ║
╠══════════════════════════════════════════════════════════════════════════════╣
║  CÓMO USAR:                                                                 ║
║  1. Sube este script + la base Wide IRT                                    ║
║  2. Escribe: "Ejecuta el PPTX Caracterización IRT"                        ║
║                                                                             ║
║  SLIDES:                                                                    ║
║    1. Portada                                                               ║
║    2. Antecedentes generales (KPIs + sexo + edad)                         ║
║    3. Sustancia principal (torta)                                          ║
║    4. Días consumo sustancia principal                                     ║
║    5. % Consumidores + Días promedio por sustancia                        ║
║    6. Urgencias + Accidentes + Salud                                      ║
║    7. Transgresión total + por tipo                                        ║
║    8. Relaciones interpersonales + Satisfacción de vida                   ║
╚══════════════════════════════════════════════════════════════════════════════╝
"""

import glob, os, unicodedata

def _norm(s):
    return unicodedata.normalize('NFD', str(s).lower()).encode('ascii','ignore').decode()

# ── Detección de país ─────────────────────────────────────────────────────────
_PAISES = {
    'republica_dominicana':'República Dominicana','repdomini':'República Dominicana',
    'dominicana':'República Dominicana','honduras':'Honduras',
    'panama':'Panamá','panam':'Panamá','el_salvador':'El Salvador',
    'salvador':'El Salvador','mexico':'México','mexic':'México',
    'ecuador':'Ecuador','peru':'Perú','argentina':'Argentina',
    'colombia':'Colombia','chile':'Chile','bolivia':'Bolivia',
    'paraguay':'Paraguay','uruguay':'Uruguay','venezuela':'Venezuela',
    'guatemala':'Guatemala','costa_rica':'Costa Rica',
    'costarica':'Costa Rica','nicaragua':'Nicaragua',
}
def _extraer_pais(filename):
    fn = _norm(str(filename).replace('.','_'))
    for key, nombre in _PAISES.items():
        if key in fn: return nombre
    return None

def _detectar_pais(wide_file):
    import pandas as _pd
    try:
        rs = _pd.read_excel(wide_file, sheet_name='Resumen', header=None)
        for _, row in rs.iterrows():
            for v in row.tolist():
                p = _extraer_pais(str(v))
                if p: return p
    except: pass
    return _extraer_pais(os.path.basename(wide_file))

def auto_archivo_wide():
    candidatos = (
        glob.glob('/mnt/user-data/uploads/*IRT*Wide*.xlsx') +
        glob.glob('/mnt/user-data/uploads/*Wide*IRT*.xlsx') +
        glob.glob('/mnt/user-data/uploads/IRT_Base*.xlsx') +
        glob.glob('/mnt/user-data/outputs/IRT_Base_Wide*.xlsx') +
        glob.glob('/home/claude/IRT_Base_Wide.xlsx'))
    if not candidatos:
        raise FileNotFoundError('⚠  No se encontró la base Wide IRT.')
    uploads = [f for f in candidatos if 'uploads' in f]
    elegido = uploads[0] if uploads else max(candidatos, key=os.path.getsize)
    print(f'  → Base Wide: {os.path.basename(elegido)}')
    return elegido

# ══════════════════════════════════════════════════════════════════════════════
print('=' * 60)
print('  SCRIPT_IRT_Universal_PPTX_Caracterizacion  v1.1')
print('=' * 60)

INPUT_FILE  = auto_archivo_wide()
OUTPUT_FILE = '/home/claude/IRT_Presentacion_Caracterizacion.pptx'

# ── FILTRO OPCIONAL POR CENTRO ────────────────────────────────────────────────
FILTRO_CENTRO = None
# ─────────────────────────────────────────────────────────────────────────────

import pandas as pd, numpy as np, json, subprocess, sys, warnings
warnings.filterwarnings('ignore')

df = pd.read_excel(INPUT_FILE, sheet_name='Base Wide', header=1)
df.columns = [str(c) for c in df.columns]
cols = df.columns.tolist()

# Filtro de centro
_col_centro = next((c for c in cols if any(x in _norm(c) for x in
                    ['codigo del centro', 'servicio de tratamiento',
                     'centro/ servicio', 'codigo centro'])), None)
if FILTRO_CENTRO and _col_centro:
    n_antes = len(df)
    df = df[df[_col_centro].astype(str).str.strip() == FILTRO_CENTRO].copy()
    df = df.reset_index(drop=True)
    print(f'\n⚑  Filtro: "{FILTRO_CENTRO}"  ({n_antes} → {len(df)} pacientes)')
    OUTPUT_FILE = f'/home/claude/IRT_Presentacion_Caracterizacion_{FILTRO_CENTRO}.pptx'

# Solo IRT1
mask1 = df['Tiene_IRT1'] == 'Sí' if 'Tiene_IRT1' in cols else pd.Series([True]*len(df))
df1   = df[mask1].copy().reset_index(drop=True)
N     = len(df1)

print(f'\n→ {N} pacientes IRT1')

# País / servicio / período
PAIS = _detectar_pais(INPUT_FILE)
if FILTRO_CENTRO:
    SERVICIO = f'{PAIS}  —  Centro {FILTRO_CENTRO}' if PAIS else f'Centro {FILTRO_CENTRO}'
else:
    SERVICIO = PAIS if PAIS else 'Servicio de Tratamiento'

MESES = {1:'Ene',2:'Feb',3:'Mar',4:'Abr',5:'May',6:'Jun',
         7:'Jul',8:'Ago',9:'Sep',10:'Oct',11:'Nov',12:'Dic'}
hoy = pd.Timestamp.now()
fecha_col = next((c for c in cols if 'fecha de administracion' in _norm(c)), None)
PERIODO = '2025'
if fecha_col:
    fch = pd.to_datetime(df1[fecha_col], errors='coerce').dropna()
    fch = fch[(fch.dt.year >= 2010) & (fch.dt.year <= hoy.year+1)]
    if len(fch):
        f0, f1_ = fch.min(), fch.max()
        PERIODO = (f'{MESES[f0.month]} {f0.year}'
                   if f0.year==f1_.year and f0.month==f1_.month
                   else f'{MESES[f0.month]}–{MESES[f1_.month]} {f0.year}'
                   if f0.year==f1_.year
                   else f'{MESES[f0.month]} {f0.year} – {MESES[f1_.month]} {f1_.year}')

print(f'  Servicio: {SERVICIO} | Período: {PERIODO}')

# ══════════════════════════════════════════════════════════════════════════════
# DETECCIÓN DE COLUMNAS (_IRT1)
# ══════════════════════════════════════════════════════════════════════════════
def col1(kws):
    for c in cols:
        if not c.endswith('_IRT1'): continue
        if all(_norm(k) in _norm(c) for k in kws): return c
    return None

COL_SEXO = next((c for c in cols if _norm(c) in ['sexo','género','genero']), None)
COL_FN   = next((c for c in cols if 'fecha de nacimiento' in _norm(c)), None)
COL_SP   = col1(['sustancia','principal'])

SUST_NOMBRES = {
    'Alcohol':['alcohol'],'Marihuana':['marihuana','cannabis'],
    'Heroína':['heroina'],'Cocaína':['cocain'],
    'Fentanilo':['fentanil'],'Inhalables':['inhalab'],
    'Metanfetamina':['metanfet','cristal'],'Crack':['crack'],
    'Pasta Base':['pasta base','pasta'],'Sedantes':['sedant','benzod'],
    'Opiáceos':['opiod','opiac'],'Tabaco':['tabaco','nicot'],
    'Otra sustancia':['otra sust'],
}
SUST_TOTAL1 = {}
for sust, kws in SUST_NOMBRES.items():
    for c in cols:
        if not c.endswith('_IRT1'): continue
        nc = _norm(c)
        if any(_norm(k) in nc for k in kws) and ('total' in nc or '(0-28)' in nc):
            SUST_TOTAL1[sust] = c; break
SUST_ACTIVAS = list(SUST_TOTAL1.keys())

COL_SPSI = col1(['salud','psicol'])
COL_SFIS = col1(['salud','fis'])
COL_URG  = next((c for c in cols if c.endswith('_IRT1') and '5)' in c and
                  any(k in c.lower() for k in ['urgencia','hospitali','emergencia'])), None)
COL_ACC  = next((c for c in cols if c.endswith('_IRT1') and '6)' in c and
                  'accidente' in c.lower()), None)

TRANS_DEF = {'Robo / Hurto':'robo','Venta de sustancias':'venta',
             'Violencia a otras personas':'violencia',
             'Violencia intrafamiliar':'intraf','Detenido / Arrestado':'detenido'}
TRANS_COLS = {n: next((c for c in cols if c.endswith('_IRT1') and kw in c.lower()), None)
              for n,kw in TRANS_DEF.items()}

REL_MAP = {'Padre':['padre'],'Madre':['madre'],'Hijos':['hijos','hijo'],
           'Hermanos':['hermanos'],'Pareja':['pareja'],'Amigos':['amigos'],'Otros':['otros']}
REL_COLS = {}
for vin, kws in REL_MAP.items():
    for c in cols:
        if not c.endswith('_IRT1'): continue
        nc = _norm(c)
        if '14)' not in c and 'relaci' not in nc: continue
        if any(_norm(k) in nc for k in kws): REL_COLS[vin]=c; break

SAT_MAP = {
    'Vida en general':[['16)'],['satisfac','vida']],
    'Lugar donde vive':[['17)'],['satisfac','lugar']],
    'Situación laboral':[['18)'],['satisfac','labor','educac']],
    'Tiempo libre':[['19)'],['satisfac','tiempo']],
    'Cap. económica':[['20)'],['satisfac','econom']],
}
SAT_COLS = {}
for dim,(nums,kws) in SAT_MAP.items():
    for c in cols:
        if not c.endswith('_IRT1'): continue
        nc = _norm(c)
        if any(n in c for n in nums) and any(_norm(k) in nc for k in kws):
            SAT_COLS[dim]=c; break

# ══════════════════════════════════════════════════════════════════════════════
# CALCULAR DATOS
# ══════════════════════════════════════════════════════════════════════════════
print('\n→ Calculando datos...')

def norm_sust(s):
    if pd.isna(s) or str(s).strip() in ['0','']: return None
    s = _norm(str(s))
    if any(x in s for x in ['alcohol','cerveza','licor']): return 'Alcohol'
    if any(x in s for x in ['marihu','cannabis','marij']): return 'Marihuana'
    if any(x in s for x in ['crack','piedra','paco']):     return 'Crack'
    if any(x in s for x in ['pasta base','pasta']):        return 'Pasta Base'
    if any(x in s for x in ['cocain','perico']):           return 'Cocaína'
    if any(x in s for x in ['fentanil']):                  return 'Fentanilo'
    if any(x in s for x in ['inhalab']):                   return 'Inhalables'
    if any(x in s for x in ['metanfet','cristal']):        return 'Metanfetamina'
    if any(x in s for x in ['sedant','benzod']):           return 'Sedantes'
    if any(x in s for x in ['heroina','opiod','morfin']):  return 'Heroína'
    if any(x in s for x in ['tabaco','nicot']):            return 'Tabaco'
    return 'Otra sustancia'

def prom1(col):
    if col is None: return None, 0
    v = pd.to_numeric(df1[col], errors='coerce').dropna()
    return (round(float(v.mean()),1) if len(v) else None), len(v)

def sino1(col):
    if col is None: return None, None
    v = df1[col].dropna().astype(str).str.strip().str.lower()
    nv = len(v); nsi = int(v.isin(['sí','si']).sum())
    return nsi, round(nsi/nv*100,1) if nv else 0.0

# Sexo
R_sexo = {}
if COL_SEXO:
    R_sexo = {k:int(v) for k,v in df1[COL_SEXO].value_counts(dropna=True).items()}
n_hombre = R_sexo.get('Hombre', R_sexo.get('H', max(R_sexo.values()) if R_sexo else 0))
n_mujer  = R_sexo.get('Mujer',  R_sexo.get('M', min(R_sexo.values()) if len(R_sexo)>1 else 0))
N_sx     = n_hombre + n_mujer
pct_h    = round(n_hombre/N_sx*100,1) if N_sx else 0

# Edad
edad_media = 0; edad_grupos = []
if COL_FN:
    edades = ((hoy - pd.to_datetime(df1[COL_FN], errors='coerce')).dt.days/365.25).dropna()
    edades = edades[(edades>=10)&(edades<=100)]
    if len(edades):
        edad_media = round(float(edades.mean()),1)
        bins = [0,17,25,35,45,55,200]
        labs = ['<18','18–25','26–35','36–45','46–55','56+']
        ec = pd.cut(edades,bins=bins,labels=labs)
        total_e = len(edades)
        edad_grupos = [{'label':l,'pct':round(int((ec==l).sum())/total_e*100,1)}
                       for l in labs if int((ec==l).sum())>0]

# Sustancia principal
sust_pp = df1[COL_SP].apply(norm_sust).dropna() if COL_SP else pd.Series(dtype=str)
R_sp = sust_pp.value_counts().to_dict() if len(sust_pp) else {}
N_sp = len(sust_pp)
sust_top1     = sust_pp.value_counts().index[0] if N_sp else '—'
sust_top1_pct = round(sust_pp.value_counts().iloc[0]/N_sp*100,1) if N_sp else 0

# Días consumo por sustancia PRINCIPAL (solo quienes la declaran como principal)
dias_pp = []
for sust, col in SUST_TOTAL1.items():
    msk = sust_pp == sust
    vals = pd.to_numeric(df1.loc[msk.reindex(df1.index,fill_value=False), col],
                         errors='coerce').dropna()
    vals = vals[vals>0]
    if len(vals): dias_pp.append({'label':sust,'prom':round(float(vals.mean()),1),'n':int(len(vals))})
dias_pp.sort(key=lambda x:-x['prom'])

# % consumidores
cons_pct = []
for sust, col in SUST_TOTAL1.items():
    v = pd.to_numeric(df1[col], errors='coerce').fillna(0)
    n_c = int((v>0).sum())
    if n_c: cons_pct.append({'label':sust,'pct':round(n_c/N*100,1),'n':n_c})
cons_pct.sort(key=lambda x:-x['pct'])

# Días promedio por sustancia (todos los consumidores)
dias_sust = []
for sust, col in SUST_TOTAL1.items():
    v = pd.to_numeric(df1[col], errors='coerce'); sub = v[v>0].dropna()
    if len(sub): dias_sust.append({'label':sust,'prom':round(float(sub.mean()),1),'n':int(len(sub))})
dias_sust.sort(key=lambda x:-x['prom'])

# Urgencias y accidentes
n_urg, pct_urg = sino1(COL_URG)
n_acc, pct_acc = sino1(COL_ACC)

# Salud
R_salud = []
for nombre, col in [('Salud Psicológica (0–10)', COL_SPSI),
                    ('Salud Física (0–10)',       COL_SFIS)]:
    m, nv = prom1(col)
    if m is not None: R_salud.append({'label':nombre,'prom':m})

# Transgresión
# Detectar formato: Sí/No dicotómico  vs  numérico (nº de veces)
def _trans_positivo(series):
    """Devuelve Serie booleana: True si el paciente cometió la transgresión.
       Soporta formato Sí/No Y formato numérico (valor > 0)."""
    s = series.dropna()
    muestra = s.astype(str).str.lower().head(20)
    es_dicotomica = muestra.isin(['sí','si','no','no aplica']).any()
    if es_dicotomica:
        return series.astype(str).str.lower().isin(['sí','si'])
    else:
        # Formato numérico: reemplazar texto ('nunca','no','no aplica') por 0
        num = pd.to_numeric(
            series.astype(str).str.lower()
                  .str.replace('nunca','0').str.replace('no aplica','0')
                  .str.replace('no','0').str.strip(),
            errors='coerce').fillna(0)
        return num > 0

any_tr_series = []
for col in [c for c in TRANS_COLS.values() if c]:
    any_tr_series.append(_trans_positivo(df1[col]))
if any_tr_series:
    any_tr = pd.concat(any_tr_series, axis=1).any(axis=1)
else:
    any_tr = pd.Series(False, index=df1.index)
n_tr = int(any_tr.sum()); pct_tr = round(n_tr/N*100,1)

trans_tipos = []
for nombre, col in TRANS_COLS.items():
    if col is None: continue
    mask = _trans_positivo(df1[col])
    nsi = int(mask.sum())
    if nsi: trans_tipos.append({'label':nombre,'n':nsi,'pct':round(nsi/N*100,1)})

# Relaciones interpersonales (% positivas)
R_rel = []
for vin, col in REL_COLS.items():
    vals = df1[col].dropna()
    vals = vals[vals.astype(str).str.lower() != 'no aplica']
    nv = len(vals)
    if nv == 0: continue
    pos = round(int(vals.astype(str).isin(['Excelente','Buena']).sum())/nv*100,1)
    R_rel.append({'label':vin,'pos':pos,'n':nv})

# Satisfacción de vida
R_sat = []
for dim, col in SAT_COLS.items():
    m, nv = prom1(col)
    if m is not None: R_sat.append({'label':dim,'prom':m})

print(f'  ✓ Calculados | Sust: {SUST_ACTIVAS}')
print(f'  Trans: {pct_tr}% | Urgencias: {pct_urg}% | Accidentes: {pct_acc}%')

# ══════════════════════════════════════════════════════════════════════════════
# JSON
# ══════════════════════════════════════════════════════════════════════════════
data = {
    'meta': {
        'servicio':       SERVICIO,
        'periodo':        PERIODO,
        'N':              N,
        'pct_h':          pct_h,
        'n_h':            n_hombre,
        'edad_media':     edad_media,
        'sust_top1':      sust_top1,
        'sust_top1_pct':  sust_top1_pct,
    },
    'sexo':      [{'label':k,'n':v} for k,v in R_sexo.items()],
    'edad':      edad_grupos,
    'sust':      [{'label':k,'pct':round(v/N_sp*100,1),'n':v} for k,v in R_sp.items()],
    'dias_pp':   dias_pp,
    'consumo':   cons_pct,
    'dias':      dias_sust,
    'urgencias': {'n': n_urg or 0, 'pct': pct_urg or 0},
    'accidentes':{'n': n_acc or 0, 'pct': pct_acc or 0},
    'salud':     R_salud,
    'transTotal':{'n': n_tr, 'pct': pct_tr},
    'transTipos':trans_tipos,
    'rel':       R_rel,
    'sat':       R_sat,
}


# ══════════════════════════════════════════════════════════════════════════════
# PYTHON-PPTX — Construcción del PowerPoint (sin Node.js)
# ══════════════════════════════════════════════════════════════════════════════
import io, matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

C_DARK  = RGBColor(0x1F,0x38,0x64)
C_MID   = RGBColor(0x2E,0x75,0xB6)
C_LIGHT = RGBColor(0xBD,0xD7,0xEE)
C_ACC   = RGBColor(0x00,0xB0,0xF0)
C_WHITE = RGBColor(0xFF,0xFF,0xFF)
C_GRAY  = RGBColor(0x59,0x59,0x59)
PIE_COLS = ['#2E75B6','#1F3864','#00B0F0','#9DC3E6','#70AD47',
            '#4472C4','#D9D9D9','#C00000','#ED7D31','#FFC000']

SLIDE_W = Inches(10); SLIDE_H = Inches(5.625)

def _rgb(h): return RGBColor(int(h[:2],16),int(h[2:4],16),int(h[4:],16))

def add_rect(sl, x,y,w,h,fill):
    s = sl.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb=fill; s.line.fill.background()

def add_txt(sl,text,x,y,w,h,size=11,bold=False,color=None,align=PP_ALIGN.LEFT,italic=False):
    tb = sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap=True
    p  = tf.paragraphs[0]; p.alignment=align
    r  = p.add_run(); r.text=str(text)
    r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic
    if color: r.font.color.rgb=color

def hdr(sl, txt):
    add_rect(sl,0,0,10,0.72,C_DARK)
    add_rect(sl,5.5,0,4.5,0.72,C_ACC)
    add_txt(sl,txt,0.25,0.05,9.5,0.62,size=18,bold=True,color=C_WHITE)

def ftr(sl, txt):
    add_txt(sl,txt,0.25,5.32,9.5,0.25,size=8,color=C_GRAY,align=PP_ALIGN.CENTER,italic=True)

def divv(sl,x):
    ln=sl.shapes.add_shape(1,Inches(x),Inches(0.78),Inches(0.02),Inches(4.85))
    ln.fill.solid(); ln.fill.fore_color.rgb=RGBColor(0xD9,0xD9,0xD9); ln.line.fill.background()

def fig2img(sl,fig,x,y,w,h):
    buf=io.BytesIO()
    fig.savefig(buf,format='png',dpi=130,bbox_inches='tight',facecolor='white')
    buf.seek(0); plt.close(fig)
    sl.shapes.add_picture(buf,Inches(x),Inches(y),Inches(w),Inches(h))

def axstyle(ax,horiz=False):
    ax.set_facecolor('white')
    (ax.xaxis if horiz else ax.yaxis).grid(True,color='#E2E8F0',linewidth=0.6,zorder=0)
    ax.set_axisbelow(True)
    for sp in ['top','right']: ax.spines[sp].set_visible(False)
    ax.spines['left'].set_color('#D0D0D0'); ax.spines['bottom'].set_color('#D0D0D0')

# ── Gráficos ──────────────────────────────────────────────────────────────────
def g_pie(labels,values,title=''):
    fig,ax=plt.subplots(figsize=(4.5,3.8))
    cols=PIE_COLS[:len(values)]
    wedges,_,at=ax.pie(values,labels=None,colors=cols,
        autopct=lambda p:f'{p:.1f}%' if p>4 else '',startangle=140,pctdistance=0.72,
        wedgeprops={'edgecolor':'white','linewidth':1.5})
    for a in at: a.set_fontsize(9); a.set_color('white'); a.set_fontweight('bold')
    ax.legend(wedges,[f'{l} (n={v})' for l,v in zip(labels,values)],
              loc='lower center',bbox_to_anchor=(0.5,-0.18),ncol=2,fontsize=8,frameon=False)
    ax.set_aspect('equal'); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

def g_bar_v(labels,values,color='#1F3864',ylabel='',max_val=None):
    fig,ax=plt.subplots(figsize=(max(4.5,len(labels)*0.85),3.5))
    bars=ax.bar(labels,values,color=color,width=0.55,zorder=3)
    for b,v in zip(bars,values):
        ax.text(b.get_x()+b.get_width()/2,b.get_height()+0.3,
                str(v),ha='center',va='bottom',fontsize=9,fontweight='bold',color='#333')
    if max_val: ax.set_ylim(0,max_val*1.15)
    else: ax.set_ylim(0,max(values)*1.3 if values else 1)
    if ylabel: ax.set_ylabel(ylabel,fontsize=8,color='#595959')
    ax.set_xticklabels(labels,fontsize=9,rotation=20,ha='right')
    axstyle(ax); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

def g_bar_h(labels,values,color='#2E75B6',xlabel=''):
    fig,ax=plt.subplots(figsize=(4.5,max(2.8,len(labels)*0.55)))
    bars=ax.barh(labels,values,color=color,height=0.55,zorder=3)
    for b,v in zip(bars,values):
        ax.text(b.get_width()+0.5,b.get_y()+b.get_height()/2,
                str(v),va='center',fontsize=9,fontweight='bold',color='#333')
    ax.set_xlim(0,max(values)*1.35 if values else 1)
    if xlabel: ax.set_xlabel(xlabel,fontsize=8,color='#595959')
    axstyle(ax,horiz=True); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

def g_kpi_sexo(sexo_list, N):
    fig,ax=plt.subplots(figsize=(4.2,3.5))
    labs=[s['label'] for s in sexo_list]; vals=[s['n'] for s in sexo_list]
    cols=['#2E75B6','#9DC3E6','#BDD7EE'][:len(labs)]
    bars=ax.bar(labs,vals,color=cols,width=0.5,zorder=3)
    for b,v in zip(bars,vals):
        ax.text(b.get_x()+b.get_width()/2,b.get_height()+0.3,
                f'{v}\n({round(v/N*100,1)}%)',ha='center',va='bottom',fontsize=10,fontweight='bold',color='#333')
    ax.set_ylim(0,max(vals)*1.4 if vals else 1)
    axstyle(ax); fig.patch.set_facecolor('white'); fig.tight_layout()
    return fig

# ── Construcción slides ───────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width=SLIDE_W; prs.slide_height=SLIDE_H
blank=prs.slide_layouts[6]
TITULO = f'Caracterización al Ingreso · {SERVICIO}'
PIE_TXT = f'N = {N}  ·  {SERVICIO}  ·  {PERIODO}  ·  Instrumento IRT'

# SLIDE 1: PORTADA
sl=prs.slides.add_slide(blank)
add_rect(sl,0,0,4.0,5.625,C_DARK)
add_rect(sl,3.1,0,1.5,5.625,C_ACC)
add_txt(sl,'Caracterización',0.25,1.6,3.2,0.7,size=22,bold=True,color=C_WHITE)
add_txt(sl,'Ingreso a Tratamiento · IRT',0.25,2.35,3.2,0.55,size=12,color=C_LIGHT)
add_txt(sl,SERVICIO.upper(),0.25,3.1,3.2,0.6,size=13,bold=True,color=C_WHITE)
add_txt(sl,PERIODO,0.25,3.75,3.2,0.4,size=11,color=C_LIGHT)
add_txt(sl,f'N = {N} personas al ingreso',0.25,4.3,3.2,0.4,size=10,color=C_LIGHT)
add_txt(sl,'IRT 1\nIngreso a Tratamiento',4.6,1.55,5.1,1.4,size=28,bold=True,color=C_GRAY,align=PP_ALIGN.CENTER)
add_txt(sl,SERVICIO.upper(),4.6,3.1,5.1,0.45,size=15,bold=True,color=C_MID,align=PP_ALIGN.CENTER)
add_txt(sl,PERIODO,4.6,3.62,5.1,0.35,size=13,bold=True,color=C_MID,align=PP_ALIGN.CENTER)
add_txt(sl,f'N = {N} personas al ingreso a tratamiento',4.6,4.1,5.1,0.35,size=10,color=C_GRAY,align=PP_ALIGN.CENTER)

# SLIDE 2: ANTECEDENTES GENERALES (KPIs + Sexo + Edad)
sl=prs.slides.add_slide(blank)
hdr(sl,TITULO); divv(sl,4.95)
# KPIs
for i,(val,lab) in enumerate([(str(N),'Personas\ningresaron'),
                               (f'{pct_h}%','Son\nhombres'),
                               (str(edad_media),'Edad\npromedio')]):
    xk=0.18+i*1.55
    add_rect(sl,xk,0.82,1.42,0.88,RGBColor(0xEE,0xF4,0xFB))
    add_txt(sl,val,xk,0.86,1.42,0.48,size=20,bold=True,color=C_DARK,align=PP_ALIGN.CENTER)
    add_txt(sl,lab,xk,1.34,1.42,0.34,size=9,color=C_GRAY,align=PP_ALIGN.CENTER)
# Sexo
add_txt(sl,'Distribución por Sexo',0.25,1.85,4.5,0.35,size=12,bold=True,color=C_GRAY)
if data['sexo']:
    fig_s=g_kpi_sexo(data['sexo'],N)
    fig2img(sl,fig_s,0.25,2.2,4.5,3.15)
# Edad
add_txt(sl,'Distribución por Rango de Edad',5.15,0.85,4.6,0.35,size=12,bold=True,color=C_GRAY)
if edad_media:
    add_txt(sl,f'Promedio: {edad_media} años',5.15,1.22,4.6,0.28,size=10,color=C_GRAY,italic=True)
if data['edad']:
    labs_e=[e['label'] for e in data['edad']]; vals_e=[e['pct'] for e in data['edad']]
    fig_e=g_bar_v(labs_e,vals_e,color='#1F3864',ylabel='%',max_val=100)
    fig2img(sl,fig_e,5.1,1.52,4.65,3.85)
ftr(sl,PIE_TXT)

# SLIDE 3: SUSTANCIA PRINCIPAL
sl=prs.slides.add_slide(blank)
hdr(sl,TITULO)
add_txt(sl,'CONSUMO SUSTANCIA PRINCIPAL AL INGRESO',1.5,0.82,7.0,0.38,
        size=14,bold=True,color=C_MID,align=PP_ALIGN.CENTER)
if data['sust']:
    labs_s=[s['label'] for s in data['sust']]; vals_s=[s['pct'] for s in data['sust']]
    fig_t=g_pie(labs_s,vals_s)
    fig2img(sl,fig_t,1.5,1.2,7.0,4.0)
add_txt(sl,f'Sustancia más frecuente: {sust_top1} ({sust_top1_pct}%)  ·  N = {N}',
        1.0,5.3,8.0,0.25,size=9,color=C_GRAY,align=PP_ALIGN.CENTER,italic=True)

# SLIDE 4: DÍAS CONSUMO SUSTANCIA PRINCIPAL
sl=prs.slides.add_slide(blank)
hdr(sl,TITULO)
add_txt(sl,'PROMEDIO DE DÍAS DE CONSUMO DE LA SUSTANCIA PRINCIPAL\nÚltimas 4 semanas · solo personas cuya sust. principal corresponde a cada categoría',
        0.5,0.82,9.0,0.65,size=12,bold=True,color=C_MID,align=PP_ALIGN.CENTER)
if data['dias_pp']:
    labs_dp=[d['label'] for d in data['dias_pp']]; vals_dp=[d['prom'] for d in data['dias_pp']]
    fig_dp=g_bar_v(labs_dp,vals_dp,color='#1F3864',ylabel='Días (0–28)',max_val=28)
    fig2img(sl,fig_dp,0.8,1.55,8.4,3.85)
ftr(sl,f'N = {N}  ·  Escala: días en últimas 4 semanas (0–28)')

# SLIDE 5: % CONSUMIDORES + DÍAS POR SUSTANCIA
sl=prs.slides.add_slide(blank)
hdr(sl,TITULO); divv(sl,4.95)
add_txt(sl,'% DE PERSONAS QUE CONSUME\nCada sustancia al ingreso',0.25,0.82,4.5,0.65,size=11,bold=True,color=C_MID)
if data['consumo']:
    labs_c=[d['label'] for d in data['consumo']]; vals_c=[d['pct'] for d in data['consumo']]
    fig_c=g_bar_v(labs_c,vals_c,color='#1F3864',ylabel='%',max_val=100)
    fig2img(sl,fig_c,0.2,1.52,4.5,3.85)
add_txt(sl,'PROMEDIO DE DÍAS DE CONSUMO\nPor sustancia (solo consumidores)',5.15,0.82,4.6,0.65,size=11,bold=True,color=C_MID)
if data['dias']:
    labs_d=[d['label'] for d in data['dias']]; vals_d=[d['prom'] for d in data['dias']]
    fig_d=g_bar_v(labs_d,vals_d,color='#2E75B6',ylabel='Días (0–28)',max_val=28)
    fig2img(sl,fig_d,5.15,1.52,4.6,3.85)
ftr(sl,PIE_TXT)

# SLIDE 6: URGENCIAS + ACCIDENTES + SALUD
sl=prs.slides.add_slide(blank)
hdr(sl,TITULO)
add_txt(sl,'SALUD Y FUNCIONAMIENTO',0.25,0.82,9.5,0.35,size=13,bold=True,color=C_MID)

def kpi_box(sl,x,y,val,label,sublabel=''):
    add_rect(sl,x,y,2.8,1.1,RGBColor(0xEE,0xF4,0xFB))
    add_txt(sl,val,x,y+0.05,2.8,0.55,size=22,bold=True,color=C_DARK,align=PP_ALIGN.CENTER)
    add_txt(sl,label,x,y+0.6,2.8,0.3,size=10,color=C_MID,align=PP_ALIGN.CENTER)
    if sublabel: add_txt(sl,sublabel,x,y+0.88,2.8,0.2,size=8,color=C_GRAY,align=PP_ALIGN.CENTER,italic=True)

n_u=data['urgencias']['n']; p_u=data['urgencias']['pct']
n_a=data['accidentes']['n']; p_a=data['accidentes']['pct']
kpi_box(sl,0.25,1.3,f'{n_u} ({p_u}%)','Con urgencia/hospitalización','por consumo')
kpi_box(sl,3.35,1.3,f'{n_a} ({p_a}%)','Con accidente','relacionado al consumo')

if data['salud']:
    labs_sal=[s['label'] for s in data['salud']]; vals_sal=[s['prom'] for s in data['salud']]
    fig_sal,ax=plt.subplots(figsize=(3.5,2.5))
    bars=ax.bar(labs_sal,vals_sal,color=['#2E75B6','#00B0F0'],width=0.5,zorder=3)
    for b,v in zip(bars,vals_sal):
        ax.text(b.get_x()+b.get_width()/2,b.get_height()+0.1,str(v),
                ha='center',va='bottom',fontsize=12,fontweight='bold',color='#333')
    ax.set_ylim(0,11); ax.axhline(5,color='#BFBFBF',linestyle='--',linewidth=0.8)
    ax.set_ylabel('Promedio (0–10)',fontsize=8,color='#595959')
    axstyle(ax); fig_sal.patch.set_facecolor('white'); fig_sal.tight_layout()
    fig2img(sl,fig_sal,6.5,1.2,3.3,2.8)
add_txt(sl,'Autopercepción de Salud (0–10)',6.5,0.85,3.3,0.35,size=11,bold=True,color=C_MID)
ftr(sl,PIE_TXT)

# SLIDE 7: TRANSGRESIÓN
sl=prs.slides.add_slide(blank)
hdr(sl,TITULO); divv(sl,4.95)
n_tr_val=data['transTotal']['n']; pct_tr_val=data['transTotal']['pct']
add_txt(sl,'TRANSGRESIÓN A LA NORMA SOCIAL',0.25,0.82,4.5,0.35,size=12,bold=True,color=C_GRAY)
# Torta trans
fig_tr,ax=plt.subplots(figsize=(3.5,3.2))
vals_tr=[n_tr_val,N-n_tr_val]; labs_tr=['Con transgresión','Sin transgresión']
wedges,_,at=ax.pie(vals_tr,labels=None,colors=['#C00000','#BDD7EE'],
    autopct=lambda p:f'{p:.1f}%',startangle=90,pctdistance=0.65,
    wedgeprops={'edgecolor':'white','linewidth':1.5})
for a in at: a.set_fontsize(10); a.set_color('white'); a.set_fontweight('bold')
ax.legend(wedges,[f'{l} (n={v})' for l,v in zip(labs_tr,vals_tr)],
          loc='lower center',bbox_to_anchor=(0.5,-0.12),fontsize=9,frameon=False)
ax.set_aspect('equal'); fig_tr.patch.set_facecolor('white'); fig_tr.tight_layout()
fig2img(sl,fig_tr,0.3,1.2,4.4,3.9)

add_txt(sl,'DISTRIBUCIÓN POR TIPO DE TRANSGRESIÓN',5.15,0.82,4.6,0.35,size=11,bold=True,color=C_GRAY)
if data['transTipos']:
    labs_tt=[t['label'] for t in data['transTipos']]; vals_tt=[t['pct'] for t in data['transTipos']]
    fig_tt=g_bar_h(labs_tt,vals_tt,color='#C00000',xlabel='% del total')
    fig2img(sl,fig_tt,5.1,1.2,4.65,4.0)
ftr(sl,PIE_TXT)

# SLIDE 8: RELACIONES + SATISFACCIÓN
sl=prs.slides.add_slide(blank)
hdr(sl,TITULO); divv(sl,4.95)

add_txt(sl,'RELACIONES INTERPERSONALES\n(% Excelente + Buena)',0.25,0.82,4.5,0.55,size=11,bold=True,color=C_MID)
if data['rel']:
    labs_r=[r['label'] for r in data['rel']]; vals_r=[r['pos'] for r in data['rel']]
    fig_r=g_bar_h(labs_r,vals_r,color='#2E75B6',xlabel='% positivo')
    fig2img(sl,fig_r,0.2,1.45,4.5,3.85)

add_txt(sl,'SATISFACCIÓN DE VIDA\nPromedio por dimensión (0–10)',5.15,0.82,4.6,0.55,size=11,bold=True,color=C_MID)
if data['sat']:
    labs_sat=[s['label'] for s in data['sat']]; vals_sat=[s['prom'] for s in data['sat']]
    fig_sat=g_bar_h(labs_sat,vals_sat,color='#00B0F0',xlabel='Promedio (0–10)')
    fig2img(sl,fig_sat,5.15,1.45,4.6,3.85)
ftr(sl,PIE_TXT)

prs.save(OUTPUT_FILE)
print(f'\n{"="*60}')
print(f'  ✅  LISTO  →  {OUTPUT_FILE}')
print(f'      {N} pacientes IRT1  ·  {PERIODO}')
print(f'{"="*60}')
